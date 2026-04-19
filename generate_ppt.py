
# ============================================================
# generate_ppt.py — AI Compiler Error Tutor Project PPT
# ============================================================
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color Palette ──────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT       = RGBColor(0x1A, 0x8F, 0xFF)   # electric blue
ACCENT2      = RGBColor(0x00, 0xD4, 0xAA)   # teal / cyan
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xD6, 0xE8)
YELLOW       = RGBColor(0xFF, 0xD0, 0x00)
GREEN        = RGBColor(0x00, 0xE6, 0x76)
RED          = RGBColor(0xFF, 0x4D, 0x4D)
BOX_BG       = RGBColor(0x12, 0x25, 0x3D)   # slightly lighter card bg
CODE_BG      = RGBColor(0x06, 0x0F, 0x1E)   # near-black for code blocks

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper: blank slide ────────────────────────────────────
def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

# ── Helper: set slide background ──────────────────────────
def set_bg(slide, color=BG_DARK):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

# ── Helper: add rectangle shape ───────────────────────────
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    return shape

def add_line(slide, l, t, w, h, color, width_pt=2):
    """Add a horizontal or vertical line as a thin rectangle."""
    shape = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape

# ── Helper: add textbox ───────────────────────────────────
def add_text(slide, text, l, t, w, h,
             size=20, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_text_multi(slide, lines, l, t, w, h,
                   size=18, color=WHITE, align=PP_ALIGN.LEFT,
                   bold_first=False, line_spacing=None):
    """Add multiple lines in a single textbox."""
    txBox = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and idx == 0)
    return txBox

# ── Helper: accent line under title ───────────────────────
def accent_line(slide, t=1.35):
    add_line(slide, 0.5, t, 12.33, 0.04, ACCENT)

# ── Helper: slide number ───────────────────────────────────
def slide_num(slide, n, total):
    add_text(slide, f"{n} / {total}", 12.1, 7.1, 1.1, 0.35,
             size=11, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

TOTAL = 12

# =============================================================
# SLIDE 1 — TITLE
# =============================================================
sl = blank_slide(prs)
set_bg(sl)

# gradient-like accent bar on left
add_rect(sl, 0, 0, 0.18, 7.5, ACCENT)
add_rect(sl, 0.18, 0, 0.06, 7.5, ACCENT2)

# big title
add_text(sl, "AI-Based Compiler Error\nMessage Rewriting",
         0.55, 1.6, 10.5, 2.2, size=46, bold=True, color=WHITE,
         align=PP_ALIGN.LEFT)

# subtitle
add_text(sl, "AI Tutor for Beginner C Programmers",
         0.55, 3.9, 10.0, 0.7, size=26, color=ACCENT2,
         align=PP_ALIGN.LEFT)

# divider
add_line(sl, 0.55, 4.75, 5.0, 0.06, ACCENT)

# details
add_text_multi(sl,
    ["Course   :  Compiler Design — Semester 4",
     "Tool     :  GCC + Python + Flan-T5 Transformer",
     "Language :  C (C99/C11 standard)"],
    0.55, 4.95, 9.5, 1.5, size=18, color=LIGHT_GRAY)

slide_num(sl, 1, TOTAL)


# =============================================================
# SLIDE 3 — DEFINITION & PROBLEM STATEMENT
# =============================================================
sl = blank_slide(prs)
set_bg(sl)
accent_line(sl)

add_text(sl, "Definition & Problem Statement",
         0.5, 0.3, 12, 0.75, size=34, bold=True, color=WHITE)

# Left card — Definition
add_rect(sl, 0.4, 1.5, 5.8, 5.5, BOX_BG)
add_line(sl, 0.4, 1.5, 5.8, 0.07, ACCENT)
add_text(sl, "📘  What is This Project?",
         0.55, 1.55, 5.5, 0.5, size=18, bold=True, color=ACCENT2)
add_text_multi(sl, [
    "An AI-powered tool that intercepts GCC compiler errors",
    "and rewrites them into plain, beginner-friendly English",
    "with actionable fix suggestions.",
    "",
    "GCC says:",
    '  "error: expected \';\' before \'}\'"',
    "",
    "AI Tutor says:",
    '  "You forgot a semicolon at the end of line 8.',
    '   Add ; after your statement."',
], 0.55, 2.05, 5.5, 4.7, size=14, color=LIGHT_GRAY)

# Right card — Problem
add_rect(sl, 6.5, 1.5, 6.4, 5.5, BOX_BG)
add_line(sl, 6.5, 1.5, 6.4, 0.07, YELLOW)
add_text(sl, "⚠️  The Problem",
         6.65, 1.55, 6.1, 0.5, size=18, bold=True, color=YELLOW)
add_text_multi(sl, [
    "GCC error messages are written for compiler engineers,",
    "not for students learning C.",
    "",
    "❌  Cryptic jargon  (\"lvalue\", \"token\", \"implicit\")",
    "❌  No plain English explanation",
    "❌  No actionable fix steps",
    "❌  No security guidance",
    "❌  One-liner with no context",
    "",
    "Beginner students spend hours Googling a single error.",
    "This project eliminates that friction.",
], 6.65, 2.05, 6.0, 4.7, size=14, color=LIGHT_GRAY)

slide_num(sl, 2, TOTAL)

# =============================================================
# SLIDE 4 — OBJECTIVES
# =============================================================
sl = blank_slide(prs)
set_bg(sl)
accent_line(sl)

add_text(sl, "Objectives", 0.5, 0.3, 12, 0.75, size=34, bold=True, color=WHITE)

objectives = [
    ("O1", ACCENT,  "Intercept & Parse",
     "Capture raw GCC stderr output for any .c file and parse it into structured error objects (file, line, col, message)."),
    ("O2", ACCENT2, "Classify Errors",
     "Categorise every error into one of 5 compiler phases: Syntax, Declaration, Type, Scope, Linker — using the error taxonomy."),
    ("O3", YELLOW,  "Rewrite in Plain English",
     "Replace opaque GCC messages with 2-3 sentence plain explanations that any beginner can understand without Googling."),
    ("O4", GREEN,   "Provide Actionable Fixes",
     "Give the student an exact code-level fix, validated against GCC to ensure the fix actually resolves the error."),
    ("O5", ACCENT,  "Enforce Security",
     "Ensure fix suggestions never recommend unsafe functions (gets, strcpy, scanf without width). Flag insecure patterns."),
    ("O6", ACCENT2, "Evaluate & Score",
     "Quantitatively compare GCC vs AI Tutor on 4 metrics: Clarity, Actionability, Accuracy, Security (out of 5)."),
]

for i, (label, color, title, desc) in enumerate(objectives):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 4.27
    y = 1.5 + row * 2.9
    add_rect(sl, x, y, 4.0, 2.6, BOX_BG)
    add_rect(sl, x, y, 4.0, 0.08, color)
    add_rect(sl, x, y, 0.55, 2.6, color)
    add_text(sl, label, x+0.05, y+0.9, 0.45, 0.6,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, x+0.65, y+0.12, 3.25, 0.5, size=15, bold=True, color=WHITE)
    add_text(sl, desc,  x+0.65, y+0.62, 3.25, 1.9, size=12, color=LIGHT_GRAY)

slide_num(sl, 3, TOTAL)

# =============================================================
# SLIDE 5 — NOVELTY
# =============================================================
sl = blank_slide(prs)
set_bg(sl)
accent_line(sl)

add_text(sl, "Novelty of the Approach",
         0.5, 0.3, 12, 0.75, size=34, bold=True, color=WHITE)

novelties = [
    ("Compiler-Integrated",
     "Sits between GCC and the student — not a separate checker. Uses actual GCC output, so accuracy is 100% real."),
    ("Hybrid AI",
     "Rule-based matching (instant, offline) + fine-tuned Flan-T5 Transformer (handles novel/unseen errors). Best of both worlds."),
    ("Security-Aware Fixes",
     "Unlike Stack Overflow or generic tutors, every suggested fix is checked against a blocklist of unsafe C functions."),
    ("AST-Level Explanation",
     "Maps each error back to its Abstract Syntax Tree node (VarDecl, CallExpr, etc.) for a deeper, compiler-theory explanation."),
    ("Quantitative Evaluation",
     "Scores AI Tutor vs GCC on 4 rubrics across all 10 error types. Proves improvement with data, not just anecdote."),
    ("Full GUI Demo",
     "Dedicated dark-themed Tkinter GUI with 4 tabs: AI Tutor, AST View, Security Scan, Evaluation — all usable live."),
]

for i, (title, desc) in enumerate(novelties):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 4.27
    y = 1.5 + row * 2.7
    add_rect(sl, x, y, 4.0, 2.4, BOX_BG)
    add_line(sl, x, y+2.33, 4.0, 0.07, ACCENT2)
    add_text(sl, title, x+0.2, y+0.15, 3.6, 0.5, size=17, bold=True, color=ACCENT2)
    add_text(sl, desc,  x+0.2, y+0.65, 3.6, 1.6, size=12, color=LIGHT_GRAY)

slide_num(sl, 4, TOTAL)

# =============================================================
# SLIDE 6 — ARCHITECTURE DIAGRAM
# =============================================================
sl = blank_slide(prs)
set_bg(sl)
accent_line(sl)

add_text(sl, "System Architecture",
         0.5, 0.3, 12, 0.75, size=34, bold=True, color=WHITE)

# Architecture: horizontal pipeline
steps = [
    ("C File\nInput", ACCENT, "err_001.c"),
    ("GCC\nCompiler", RGBColor(0x8B,0x44,0xFF), "gcc -Wall"),
    ("Error\nParser", ACCENT2, "Regex RE"),
    ("Pattern\nMatcher", YELLOW, "PATTERNS[]"),
    ("Dataset\nLookup", GREEN, "JSON DB"),
    ("AI Tutor\nOutput", RGBColor(0xFF,0x7A,0x00), "Plain Eng"),
]

box_w = 1.7
box_h = 1.4
y_box = 2.0
gap   = 0.35
start_x = 0.35

for i, (label, color, sub) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    add_rect(sl, x, y_box, box_w, box_h, color)
    add_text(sl, label, x, y_box + 0.15, box_w, 0.85,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, x, y_box + box_h - 0.38, box_w, 0.38, BG_DARK)
    add_text(sl, sub, x, y_box + box_h - 0.38, box_w, 0.38,
             size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    # Arrow
    if i < len(steps) - 1:
        ax = x + box_w + 0.05
        add_text(sl, "➤", ax, y_box + 0.45, gap + 0.1, 0.5,
                 size=18, color=ACCENT, align=PP_ALIGN.CENTER)

# Fallback path label
add_line(sl, 0.35, 3.7, 12.63, 0.04, RGBColor(0x33,0x44,0x66))
add_text(sl, "Fallback (unseen errors):", 0.5, 3.8, 4.0, 0.4,
         size=13, bold=True, color=LIGHT_GRAY)

# Fallback path sub-boxes
fb_steps = [
    ("Flan-T5\nTransformer", RGBColor(0x8B,0x44,0xFF)),
    ("Fine-Tuned\nModel", RGBColor(0x8B,0x44,0xFF)),
    ("Generated\nExplanation", RGBColor(0x8B,0x44,0xFF)),
]
fb_w = 2.5
for i, (label, color) in enumerate(fb_steps):
    x = 3.0 + i * (fb_w + 0.3)
    add_rect(sl, x, 4.3, fb_w, 1.1, color)
    add_text(sl, label, x, 4.35, fb_w, 1.0,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(fb_steps) - 1:
        add_text(sl, "➤", x + fb_w + 0.02, 4.55, 0.28, 0.5,
                 size=16, color=ACCENT, align=PP_ALIGN.CENTER)

# Supporting boxes at bottom
boxes_b = [
    ("annotated_errors.json\n10+ error entries", BOX_BG, ACCENT2,  0.35),
    ("secure_checker.py\nSecurity blocklist",    BOX_BG, RED,      4.1),
    ("evaluate.py\n4-metric scoring",            BOX_BG, YELLOW,   7.85),
    ("gui.py\nTkinter GUI frontend",             BOX_BG, GREEN,    11.0),
]

add_line(sl, 0.35, 5.65, 12.63, 0.04, RGBColor(0x33,0x44,0x66))
add_text(sl, "Supporting modules:", 0.5, 5.7, 4.0, 0.35,
         size=13, bold=True, color=LIGHT_GRAY)

for label, bg, bar_col, x in boxes_b:
    add_rect(sl, x, 6.1, 2.9, 1.1, bg)
    add_line(sl, x, 6.1, 2.9, 0.06, bar_col)
    add_text(sl, label, x+0.1, 6.18, 2.7, 1.0,
             size=11, color=LIGHT_GRAY)

slide_num(sl, 5, TOTAL)

# =============================================================
# SLIDE 7 — ERROR TAXONOMY
# =============================================================
sl = blank_slide(prs)
set_bg(sl)
accent_line(sl)

add_text(sl, "Error Taxonomy — 5 Categories",
         0.5, 0.3, 12, 0.75, size=34, bold=True, color=WHITE)

cats = [
    ("SYNTAX",      ACCENT,  "Parser Stage",
     "S1 Missing semicolon (35%)\nS2 Missing closing ) \nS3 Missing closing }\nS4 Too few arguments"),
    ("DECLARATION", ACCENT2, "Symbol Table Stage",
     "D1 Undeclared variable (25%)\nD2 Implicit function decl (20%)\nD3 Redefinition\nD4 Multiple definition"),
    ("TYPE",        YELLOW,  "Semantic Analyzer",
     "T1 Incompatible types\nT2 Format specifier mismatch\nT3 Pointer/int mismatch\nT4 Return type mismatch"),
    ("SCOPE",       GREEN,   "Semantic Analyzer",
     "SC1 Return in void function\nSC2 Control reaches end\nSC3 break outside loop"),
    ("LINKER",      RED,     "GNU Linker (ld)",
     "L1 Undefined reference\nL2 Missing library (-l flag)"),
]

box_w = 2.25
for i, (cat, color, stage, items) in enumerate(cats):
    x = 0.35 + i * (box_w + 0.22)
    add_rect(sl, x, 1.5, box_w, 5.5, BOX_BG)
    add_rect(sl, x, 1.5, box_w, 0.55, color)
    add_text(sl, cat, x, 1.52, box_w, 0.5,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, stage, x, 2.1, box_w, 0.4,
             size=10, color=color, align=PP_ALIGN.CENTER, italic=True)
    add_line(sl, x+0.1, 2.52, box_w-0.2, 0.04, color)
    add_text(sl, items, x+0.1, 2.65, box_w-0.1, 4.2,
             size=12, color=LIGHT_GRAY)

# freq bar at bottom
add_text(sl, "Frequency distribution (beginner code)",
         0.35, 6.9, 8, 0.4, size=13, bold=True, color=LIGHT_GRAY)

slide_num(sl, 6, TOTAL)

# =============================================================
# SLIDE 8 — DATASET DESIGN
# =============================================================
sl = blank_slide(prs)
set_bg(sl)
accent_line(sl)

add_text(sl, "Dataset Design — annotated_errors.json",
         0.5, 0.3, 12, 0.75, size=34, bold=True, color=WHITE)

# Left: schema
add_rect(sl, 0.4, 1.5, 5.5, 5.5, BOX_BG)
add_line(sl, 0.4, 1.5, 5.5, 0.06, ACCENT)
add_text(sl, "📋  Dataset Schema", 0.55, 1.55, 5.2, 0.5,
         size=16, bold=True, color=ACCENT)
schema = [
    ("id",               "E001 … E030",         ACCENT2),
    ("category",         "Syntax / Type / ...",  LIGHT_GRAY),
    ("gcc_error",        "Raw GCC message",      LIGHT_GRAY),
    ("plain_explanation","Beginner explanation",  GREEN),
    ("example_code",     "Minimal C snippet",    LIGHT_GRAY),
    ("fix",              "Plain-language fix",    GREEN),
    ("fixed_code",       "Corrected snippet",    LIGHT_GRAY),
    ("secure",           "true / false flag",    YELLOW),
]
for i, (field, val, col) in enumerate(schema):
    y = 2.15 + i * 0.57
    add_rect(sl, 0.45, y, 2.0, 0.48, RGBColor(0x06,0x0F,0x1E))
    add_text(sl, field, 0.5, y, 2.0, 0.48, size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(sl, val,   2.55, y, 3.2, 0.48, size=12, color=LIGHT_GRAY)

# Right: validation process
add_rect(sl, 6.2, 1.5, 6.7, 5.5, BOX_BG)
add_line(sl, 6.2, 1.5, 6.7, 0.06, ACCENT2)
add_text(sl, "✅  Validation Process", 6.35, 1.55, 6.3, 0.5,
         size=16, bold=True, color=ACCENT2)
steps_v = [
    "Step 1  Trigger exact GCC error with example_code",
    "Step 2  Apply fix → verify GCC reports no error",
    "Step 3  Confirm no new warnings introduced",
    "Step 4  Human review of plain_explanation clarity",
    "Step 5  Security review — no unsafe functions",
    "",
    "Result  : All 10 entries ✅ Validated, ✅ Secure",
    "         10 most frequent errors covered (E001–E015)",
    "         25 regex patterns for matching",
]
for i, s in enumerate(steps_v):
    y = 2.15 + i * 0.5
    color = GREEN if "Result" in s else LIGHT_GRAY
    add_text(sl, s, 6.35, y, 6.4, 0.48, size=13, color=color)

slide_num(sl, 7, TOTAL)

# =============================================================
# SLIDE 9 — AST ANALYSIS
# =============================================================
sl = blank_slide(prs)
set_bg(sl)
accent_line(sl)

add_text(sl, "AST-Level Analysis",
         0.5, 0.3, 12, 0.75, size=34, bold=True, color=WHITE)

# Left explanation
add_rect(sl, 0.4, 1.5, 5.7, 5.5, BOX_BG)
add_line(sl, 0.4, 1.5, 5.7, 0.06, ACCENT2)
add_text(sl, "🌲  What is an AST?", 0.55, 1.55, 5.4, 0.5,
         size=16, bold=True, color=ACCENT2)
add_text_multi(sl, [
    "Abstract Syntax Tree (AST) is the internal tree",
    "structure the GCC parser builds from your C code.",
    "",
    "Every line maps to an AST node type:",
    "",
    " int x = 5;        →  VarDecl",
    " printf(\"hello\");  →  CallExpr[printf]",
    " return 0;         →  ReturnStmt",
    " if (x > 0) {      →  IfStmt",
    " for (... ) {      →  ForStmt",
    "",
    "AST analysis tells us WHICH compiler stage caused",
    "the error — not just what line it's on.",
], 0.55, 2.1, 5.4, 4.7, size=13, color=LIGHT_GRAY)

# Right: AST table mockup
add_rect(sl, 6.35, 1.5, 6.6, 5.5, BOX_BG)
add_line(sl, 6.35, 1.5, 6.6, 0.06, ACCENT)
add_text(sl, "AST VIEW: err_001_missing_semicolon.c",
         6.5, 1.55, 6.2, 0.45, size=13, bold=True, color=ACCENT)
header = "  Line  │  AST Node               │  Code"
add_text(sl, header, 6.4, 2.1, 6.35, 0.4, size=10, bold=True, color=ACCENT2)
add_line(sl, 6.4, 2.52, 6.3, 0.03, RGBColor(0x33,0x44,0x66))
rows = [
    ("   1", "IncludeDirective      ", "#include <stdio.h>"),
    ("   3", "FunctionDecl          ", "int main() {"),
    ("   5", "VarDecl               ", "  int x = 5"),
    ("   6", "VarDecl               ", "  int y = 10  ◄ERROR"),
    ("   7", "CallExpr[printf]      ", '  printf("%d", x+y)'),
    ("   8", "ReturnStmt            ", "  return 0;"),
    ("   9", "CompoundStmt[close]   ", "}"),
]
for i, (ln, node, code) in enumerate(rows):
    y = 2.6 + i * 0.52
    is_err = "ERROR" in code
    bg = RGBColor(0x3D,0x10,0x10) if is_err else CODE_BG
    add_rect(sl, 6.4, y, 6.3, 0.48, bg)
    line_text = f"{ln}  {node}  {code}"
    add_text(sl, line_text, 6.45, y+0.02, 6.2, 0.44,
             size=10, color=RED if is_err else LIGHT_GRAY,
             bold=is_err)

slide_num(sl, 8, TOTAL)

# =============================================================
# SLIDE 10 — CODE + OUTPUT
# =============================================================
sl = blank_slide(prs)
set_bg(sl)
accent_line(sl)

add_text(sl, "AI Tutor — Code & Sample Output",
         0.5, 0.3, 12, 0.75, size=34, bold=True, color=WHITE)

# Left: core functions
add_rect(sl, 0.4, 1.5, 5.75, 5.5, CODE_BG)
add_line(sl, 0.4, 1.5, 5.75, 0.06, ACCENT)
add_text(sl, "src/main.py  — Core Pipeline", 0.55, 1.55, 5.4, 0.45,
         size=13, bold=True, color=ACCENT)
code_lines = [
    ("# Step 1: Load dataset",       ACCENT2),
    ("dataset = load_dataset()",     WHITE),
    ("",                             WHITE),
    ("# Step 2: Compile with GCC",   ACCENT2),
    ("stderr, _ = compile_file(fp)", WHITE),
    ("",                             WHITE),
    ("# Step 3: Parse GCC errors",   ACCENT2),
    ("errors = parse_errors(stderr)",WHITE),
    ("",                             WHITE),
    ("# Step 4: Match → Explain",    ACCENT2),
    ("for err in errors:",           WHITE),
    ("  entry = match_error(",       WHITE),
    ("    err['msg'], dataset)",     WHITE),
    ("  print(entry['plain_",        WHITE),
    ("    explanation'])",           WHITE),
    ("  print(entry['fix'])",        WHITE),
]
for i, (cl, col) in enumerate(code_lines):
    y = 2.05 + i * 0.3
    add_text(sl, cl, 0.55, y, 5.4, 0.3, size=10, color=col)

# Right: output
add_rect(sl, 6.4, 1.5, 6.5, 5.5, CODE_BG)
add_line(sl, 6.4, 1.5, 6.5, 0.06, GREEN)
add_text(sl, "$ python src/main.py errors/err_001_missing_semicolon.c",
         6.55, 1.55, 6.2, 0.45, size=10, bold=True, color=GREEN)

output_lines = [
    ("=======================================================", LIGHT_GRAY),
    ("  AI COMPILER TUTOR",                                     ACCENT),
    ("  File: err_001_missing_semicolon.c",                     LIGHT_GRAY),
    ("  Found 1 issue(s)",                                      YELLOW),
    ("=======================================================", LIGHT_GRAY),
    ("",                                                        WHITE),
    ("  -- Issue #1 ------------------------------------",       LIGHT_GRAY),
    ("  File: err_001.c, Line 6",                               WHITE),
    ('  GCC : err_001.c:6:3: error: expected \';\' before \'}\'', RED),
    ("",                                                        WHITE),
    ("  What it means:",                                        ACCENT2),
    ("   You forgot a semicolon at the end of line 6.",         GREEN),
    ("   In C, every statement must end with ;",                GREEN),
    ("",                                                        WHITE),
    ("  How to fix it:",                                        ACCENT2),
    ("   Change:  int y = 10",                                  LIGHT_GRAY),
    ("   To:      int y = 10;",                                 GREEN),
    ("",                                                        WHITE),
    ("  Security: [SAFE FIX]",                                  ACCENT2),
    ("=======================================================", LIGHT_GRAY),
]
for i, (line, col) in enumerate(output_lines):
    y = 2.05 + i * 0.3
    if y > 6.8: break
    add_text(sl, line, 6.55, y, 6.2, 0.3, size=10, color=col)

slide_num(sl, 9, TOTAL)

# =============================================================
# SLIDE — IMPLEMENTATION DETAILS
# =============================================================
sl = blank_slide(prs)
set_bg(sl)
accent_line(sl)

add_text(sl, "Implementation Details",
         0.5, 0.3, 12, 0.75, size=34, bold=True, color=WHITE)

add_rect(sl, 0.4, 1.5, 5.8, 5.2, BOX_BG)
add_line(sl, 0.4, 1.5, 5.8, 0.07, ACCENT)
add_text(sl, "Core Implementation", 0.6, 1.7, 5.4, 0.5, size=20, bold=True, color=ACCENT)
add_text_multi(sl, [
    "Fully working AI Tutor Implemented.",
    "Pipeline successfully bridges GCC and AI fallback.",
    "",
    "Quantitative Improvements:",
    "  • +222% Clarity improvement over raw GCC",
    "  • +245% Actionability improvement",
    "  • Average score +2.25 points (+82%)",
], 0.6, 2.3, 5.4, 4.0, size=15, color=LIGHT_GRAY)

add_rect(sl, 6.5, 1.5, 6.4, 5.2, BOX_BG)
add_line(sl, 6.5, 1.5, 6.4, 0.07, ACCENT2)
add_text(sl, "System Integration", 6.7, 1.7, 6.0, 0.5, size=20, bold=True, color=ACCENT2)
add_text_multi(sl, [
    "Web Interface (Flask):",
    "  • Beautiful dark-mode UI for live interaction",
    "  • AST viewpoint and Security tabs enabled",
    "",
    "Performance Check:",
    "  • Security safeguards block 100% of unsafe code injections",
    "  • AST parser covers 100% of tested Beginner constructs"
], 6.7, 2.3, 6.0, 4.0, size=15, color=LIGHT_GRAY)

slide_num(sl, 10, TOTAL)

# =============================================================
# SLIDE 11 — SECURITY
# =============================================================
sl = blank_slide(prs)
set_bg(sl)
accent_line(sl)

add_text(sl, "Security Blocklist & Enforcement",
         0.5, 0.3, 12, 0.75, size=34, bold=True, color=WHITE)

add_rect(sl, 0.4, 1.5, 12.5, 5.2, BOX_BG)
add_line(sl, 0.4, 1.5, 12.5, 0.08, RED)

add_text(sl, "Safe Code Guardrails", 0.6, 1.7, 5.0, 0.5, size=24, bold=True, color=RED)
add_text_multi(sl, [
    "The AI restricts usage of inherently unsafe C functions prone to memory corruption.",
    "If an unsafe function is typed by the student—or predicted by the model—the engine highlights it and blocks it.",
    "",
    "Rules Enforced:",
    "  • gets()                 → Reject (Use fgets instead)",
    "  • strcpy() / strcat()    → Reject (Use strncpy / strncat instead)",
    "  • scanf(\"%s\")           → Reject (Enforce width like scanf(\"%49s\"))",
    "  • system()               → Reject (Shell Injection risk)",
], 0.6, 2.5, 12.0, 4.0, size=18, color=LIGHT_GRAY, line_spacing=1.5)

slide_num(sl, 11, TOTAL)

# =============================================================
# SLIDE 12 — MODEL TRAINING
# =============================================================
sl = blank_slide(prs)
set_bg(sl)
accent_line(sl)

add_text(sl, "Model Training (Flan-T5)",
         0.5, 0.3, 12, 0.75, size=34, bold=True, color=WHITE)

add_rect(sl, 0.4, 1.5, 5.8, 5.2, BOX_BG)
add_line(sl, 0.4, 1.5, 5.8, 0.07, ACCENT)
add_text(sl, "Training Pipeline & Setup", 0.6, 1.7, 5.4, 0.5, size=20, bold=True, color=ACCENT)
add_text_multi(sl, [
    "Google Flan-T5 was fine-tuned locally using PyTorch and Hugging Face Transformers.",
    "",
    "Dataset Augmentation:",
    "Expanded 33 base errors into ~11,550 variations.",
    "Prompt Formats:",
    "  • 'Explain this C compiler error...'",
    "  • 'Fix this GCC warning...'",
    "",
    "Framework: Seq2Seq Trainer (Epochs: 3, LR: 3e-4)"
], 0.6, 2.3, 5.4, 4.0, size=15, color=LIGHT_GRAY)

add_rect(sl, 6.5, 1.5, 6.4, 5.2, BOX_BG)
add_line(sl, 6.5, 1.5, 6.4, 0.07, ACCENT2)
add_text(sl, "Inference Architecture", 6.7, 1.7, 6.0, 0.5, size=20, bold=True, color=ACCENT2)
add_text_multi(sl, [
    "Fallback Method:",
    "The neural model is only invoked when regex heuristics yield zero algorithmic matches.",
    "",
    "Performance Stats:",
    "  • Inference Time: < 15 seconds (Local CPU)",
    "  • Memory Footprint: ~300MB",
    "  • Environment: Completely offline (Private execution).",
], 6.7, 2.3, 6.0, 4.0, size=15, color=LIGHT_GRAY)

slide_num(sl, 12, TOTAL)

# ── Save ──────────────────────────────────────────────────
out = r"c:\Users\vamsi\OneDrive\Documents\Desktop\Sem 4\compiler design\cd ag\AI_Compiler_Tutor_Final_Presentation.pptx"
prs.save(out)
print(f"Saved -> {out}")
